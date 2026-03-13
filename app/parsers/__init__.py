import fitz  # PyMuPDF
import pdfplumber
from pathlib import Path
from typing import List

from docx import Document as DocxDocument
from pptx import Presentation

from app.models import PageContent, DocumentMetadata, ParsedDocument


class PDFParser:
    """Parser for PDF documents"""

    def __init__(self):
        pass

    def parse(self, file_path: str, language: str = "French") -> ParsedDocument:
        """
        Parse a PDF file and extract structured content.

        Args:
            file_path: Path to PDF file
            language: Document language (French or English)

        Returns:
            ParsedDocument with all extracted content
        """
        file_path = Path(file_path)

        # Extract text content using PyMuPDF
        doc = fitz.open(file_path)
        pages = []

        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text
            text = page.get_text("text")

            # Check for images
            image_list = page.get_images()
            has_images = len(image_list) > 0

            # Use visible page number (1-indexed)
            visible_page_num = self._extract_visible_page_number(page, page_num + 1)

            page_content = PageContent(
                page_number=visible_page_num,
                text=text,
                has_images=has_images,
                has_tables=False,  # Will detect tables next
                images_text=[]
            )
            pages.append(page_content)

        doc.close()

        # Enhanced table detection using pdfplumber
        pages = self._detect_tables(file_path, pages)

        # Create metadata
        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="pdf",
            total_pages=len(pages),
            language=language
        )

        return ParsedDocument(
            metadata=metadata,
            pages=pages,
            raw_metadata={}
        )

    def _extract_visible_page_number(self, page, default_num: int) -> int:
        """
        Try to extract the visible page number from the page content.
        Falls back to sequential numbering if not found.

        Args:
            page: fitz Page object
            default_num: Default page number to use

        Returns:
            Visible page number
        """
        # For V1, use default sequential numbering
        # In V2, can add heuristics to detect page numbers in headers/footers
        return default_num

    def _detect_tables(self, file_path: Path, pages: List[PageContent]) -> List[PageContent]:
        """
        Detect tables in PDF pages using pdfplumber.

        Args:
            file_path: Path to PDF
            pages: List of PageContent objects

        Returns:
            Updated list of PageContent with table detection
        """
        try:
            with pdfplumber.open(file_path) as pdf:
                for idx, pdf_page in enumerate(pdf.pages):
                    if idx < len(pages):
                        tables = pdf_page.find_tables()
                        pages[idx].has_tables = len(tables) > 0
        except Exception:
            # If table detection fails, continue without it
            pass

        return pages

    def extract_text_from_page(self, file_path: str, page_number: int) -> str:
        """
        Extract text from a specific page.

        Args:
            file_path: Path to PDF
            page_number: Page number (1-indexed)

        Returns:
            Text content from that page
        """
        doc = fitz.open(file_path)
        if page_number < 1 or page_number > len(doc):
            doc.close()
            raise ValueError(f"Page number {page_number} out of range")

        page = doc[page_number - 1]
        text = page.get_text("text")
        doc.close()

        return text


class PPTXParser:
    """Parser for PowerPoint documents"""

    def parse(self, file_path: str, language: str = "French") -> ParsedDocument:
        file_path = Path(file_path)
        presentation = Presentation(file_path)

        pages: List[PageContent] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            text_runs = []
            has_images = False
            has_tables = False

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_runs.append(shape.text)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    has_images = True
                if shape.has_table:
                    has_tables = True

            text = "\n".join([t for t in text_runs if t])
            pages.append(
                PageContent(
                    page_number=idx,
                    text=text,
                    has_images=has_images,
                    has_tables=has_tables,
                    images_text=[],
                )
            )

        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="pptx",
            total_pages=len(pages),
            language=language,
        )

        return ParsedDocument(metadata=metadata, pages=pages, raw_metadata={})


class DOCXParser:
    """Parser for Word documents"""

    def parse(self, file_path: str, language: str = "French") -> ParsedDocument:
        file_path = Path(file_path)
        doc = DocxDocument(file_path)

        # DOCX doesn't expose pages; treat the whole doc as one page for V1
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        text = "\n".join(paragraphs)
        has_tables = len(doc.tables) > 0

        pages = [
            PageContent(
                page_number=1,
                text=text,
                has_images=False,
                has_tables=has_tables,
                images_text=[],
            )
        ]

        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="docx",
            total_pages=len(pages),
            language=language,
        )

        return ParsedDocument(metadata=metadata, pages=pages, raw_metadata={})


class XLSXParser:
    """Parser for Excel workbooks — each sheet becomes one page."""

    def parse(self, file_path: str, language: str = "English") -> ParsedDocument:
        import openpyxl

        file_path = Path(file_path)
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        pages = []

        for idx, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            rows_text: List[str] = []
            for row in ws.iter_rows(values_only=True):
                cell_values = [str(c) if c is not None else "" for c in row]
                row_str = "\t".join(cell_values).strip()
                if row_str:
                    rows_text.append(row_str)

            text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text)
            pages.append(
                PageContent(
                    page_number=idx,
                    text=text,
                    has_images=False,
                    has_tables=True,
                    images_text=[],
                )
            )

        wb.close()

        metadata = DocumentMetadata(
            filename=file_path.name,
            file_type="xlsx",
            total_pages=len(pages),
            language=language,
        )

        return ParsedDocument(metadata=metadata, pages=pages, raw_metadata={})
